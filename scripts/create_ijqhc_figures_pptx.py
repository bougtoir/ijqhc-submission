#!/usr/bin/env python3
"""Create editable PowerPoint files for IJQHC figures.

Provides the figures as an editable .pptx (one figure per slide,
widescreen 13.333 x 7.5 inches) so co-authors can edit / replace figures
without regenerating the source PNGs.

English and Japanese versions are produced separately to comply with the
language-consistency rule.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'documents', 'IJQHC')
REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
FIG_DIR = os.path.join(REPO_ROOT, 'output')


def build_deck(figures, out_path, lang='en'):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, (img, title, caption) in enumerate(figures, 1):
        slide = prs.slides.add_slide(blank_layout)
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2),
                                             Inches(12.5), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(22)
        if lang == 'jp':
            run.font.name = 'MS Gothic'
        else:
            run.font.name = 'Arial'

        # Image
        img_path = os.path.join(FIG_DIR, img)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(1.8), Inches(1.0),
                                      width=Inches(9.7))
        else:
            print(f"Warning: missing image {img_path}")

        # Caption
        cap_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.4),
                                            Inches(12.5), Inches(1.0))
        cf = cap_box.text_frame
        cf.word_wrap = True
        cp = cf.paragraphs[0]
        crun = cp.add_run()
        crun.text = caption
        crun.font.size = Pt(12)
        if lang == 'jp':
            crun.font.name = 'MS Mincho'
        else:
            crun.font.name = 'Arial'

    prs.save(out_path)
    print(f"Saved: {out_path}")


# English deck
figures_en = [
    (
        'rapm_fig1_en.png',
        'Figure 1. Geographic distribution of anaesthesia standardised claim '
        'ratios across 335 secondary medical areas of Japan, fiscal year 2022.',
        '(A) General anaesthesia (L008). (B) Spinal anaesthesia (L004). '
        '(C) Epidural anaesthesia as main anaesthetic (L002). '
        '(D) Continuous epidural infusion (L003). Choropleth maps shaded '
        'by quintile of the standardised claim ratio (national average = '
        '100). Red circles mark secondary medical areas containing at '
        'least one university hospital. Areas masked by the data provider '
        'owing to low volume are shown in grey.',
    ),
    (
        'rapm_fig2_en.png',
        'Figure 2. University hospital presence and the combined '
        'general-anaesthesia plus continuous-epidural measure.',
        '(A) Distribution of secondary medical areas containing at least one '
        'university hospital (n = 64 of 335; red). (B) Choropleth map of the '
        'combined general-anaesthesia plus continuous-epidural standardised '
        'claim ratio (mean of L008 and L003 SCR; 331 areas with data for '
        'both codes), shaded by quintile. Red circles mark secondary medical '
        'areas containing at least one university hospital. Areas masked by '
        'the data provider for either code are shown in grey.',
    ),
]
build_deck(figures_en,
           os.path.join(OUTPUT_DIR, 'regional_anaesthesia_figures_IJQHC_EN.pptx'),
           lang='en')

# Japanese deck
figures_jp = [
    (
        'rapm_fig1_jp.png',
        '図1．335二次医療圏における麻酔の標準化レセプト出現比の地理分布'
        '（2022年度）。',
        '(A) 閉鎖循環式全身麻酔（L008）。(B) 脊椎麻酔（L004）。標準化'
        'レセプト出現比（全国平均=100）の五分位で陰影付け。低症例により'
        'マスクされた医療圏は灰色で示す。',
    ),
    (
        'rapm_fig2_jp.png',
        '図2．大学病院所在と全身麻酔の標準化レセプト出現比。',
        '(A) 少なくとも1つの大学病院を含む二次医療圏（335医療圏中64）の'
        '分布。(B) 47都道府県すべてにおける大学病院所在医療圏と非所在'
        '医療圏の平均L008出現比の都道府県内比較。47/47都道府県で大学'
        '病院効果が正であった（都道府県内平均差+63.3点、対応のある'
        't=13.28、P<0.001）。',
    ),
]
build_deck(figures_jp,
           os.path.join(OUTPUT_DIR, 'regional_anaesthesia_figures_IJQHC_JP.pptx'),
           lang='jp')
