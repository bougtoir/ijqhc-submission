#!/usr/bin/env python3
"""
Create JP PPTX with 1 figure/table per slide.
- Code-generated maps: embedded as images (as-is)
- Tables: editable pptx tables
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


def add_title_slide(title_text, subtitle_text=''):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.name = 'Arial'
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER
    return slide


def add_image_slide(image_path, caption, slide_title=''):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    if slide_title:
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0, 0, 0)

    if os.path.exists(image_path):
        from PIL import Image
        img = Image.open(image_path)
        img_w, img_h = img.size
        aspect = img_w / img_h

        max_w = Inches(12.0)
        max_h = Inches(5.6)

        if aspect > (12.0 / 5.6):
            w = max_w
            h = int(w / aspect)
        else:
            h = max_h
            w = int(h * aspect)

        left = int((prs.slide_width - w) / 2)
        top = Inches(0.75)

        slide.shapes.add_picture(image_path, left, top, w, h)
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f'[画像が見つかりません: {image_path}]'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 0, 0)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_table_slide(title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(12.7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0, 0, 0)

    n_rows = len(rows) + 1
    n_cols = len(headers)

    if col_widths is None:
        total_w = Inches(11.0)
        col_w = total_w / n_cols
        col_widths_emu = [int(col_w)] * n_cols
    else:
        col_widths_emu = [Inches(w) for w in col_widths]

    table_width = sum(col_widths_emu)
    table_left = int((prs.slide_width - table_width) / 2)
    table_top = Inches(1.2)
    row_height = Inches(0.4)
    table_height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    for i, w in enumerate(col_widths_emu):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.name = 'Arial'
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(51, 63, 80)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Arial'
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 244, 248)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    return slide


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════

add_title_slide(
    "日本の335二次医療圏における麻酔診療の地域差",
    "全国レセプトデータを用いた横断的分析"
)


# ═══════════════════════════════════════════════════════════
# SLIDE 2: Table 1
# ═══════════════════════════════════════════════════════════

add_table_slide(
    '表1. 分析対象の麻酔関連診療行為コード',
    ['コード', '術式名', '臨床的意義', '対象圏数'],
    [
        ['L008', '閉鎖循環式全身麻酔', 'GA量の主要指標', '334'],
        ['L002', '硬膜外麻酔', '区域麻酔（単独又は併用）', '307'],
        ['L003', '硬膜外麻酔後における局所麻酔剤の持続的注入', 'GA+硬膜外併用の直接指標', '331'],
        ['L004', '脊椎麻酔', 'GA代替法', '334'],
        ['L009', '麻酔管理料（I）', '麻酔科専門医配置の代理指標', '314'],
        ['L100', '神経ブロック（入院）', 'ペインクリニック活動指標', '335'],
    ],
    col_widths=[1.2, 4.5, 3.5, 1.8]
)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: Table 2
# ═══════════════════════════════════════════════════════════

add_table_slide(
    '表2. 二次医療圏における麻酔関連診療行為のSCR分布（全国平均=100）',
    ['コード', '術式', 'P10', 'P50', 'P90', 'CV (%)', 'n'],
    [
        ['L008', '全身麻酔', '32.3', '73.2', '131.5', '54.6', '334'],
        ['L002', '硬膜外麻酔', '7.6', '73.9', '184.0', '83.2', '307'],
        ['L003', '持続硬膜外注入', '19.0', '73.1', '150.0', '64.9', '331'],
        ['L004', '脊椎麻酔', '31.1', '84.3', '169.0', '56.3', '334'],
        ['L009', '麻酔管理料I', '32.5', '79.1', '149.1', '57.2', '314'],
        ['L100', '神経ブロック（入院）', '23.4', '72.2', '174.7', '72.1', '335'],
    ],
    col_widths=[1.0, 3.0, 1.0, 1.0, 1.0, 1.0, 1.0]
)


# ═══════════════════════════════════════════════════════════
# SLIDES 4-9: Figures 1-6 (2D maps) - JP versions
# ═══════════════════════════════════════════════════════════

figures_2d = [
    ('/home/ubuntu/map_L008_scr.png',
     '図1. 二次医療圏別の全身麻酔（L008）標準化レセプト出現比。'
     '実線=都道府県境界、破線=二次医療圏境界。北方領土は白色（医療圏未設定）。',
     '図1'),
    ('/home/ubuntu/map_univ_presence.png',
     '図2. 二次医療圏別の大学病院所在状況。赤丸は1校以上の大学病院を含むSMA（335圏中64圏）。',
     '図2'),
    ('/home/ubuntu/map_L004_scr.png',
     '図3. 二次医療圏別の脊椎麻酔（L004）標準化レセプト出現比。',
     '図3'),
    ('/home/ubuntu/map_L008_L004_combined.png',
     '図4. 全身麻酔+脊椎麻酔合計（L008+L004）SCR。合算により査定による再分類の影響を中和。',
     '図4'),
    ('/home/ubuntu/map_L003_scr.png',
     '図5. 持続硬膜外注入（L003）SCR \u2014 GA+硬膜外併用の直接指標。CV=64.9%。',
     '図5'),
    ('/home/ubuntu/map_L003_L008_ratio_corrected.png',
     '図6. L003/L008比（GA1件あたりの硬膜外併用率）。手術件数で補正済み。',
     '図6'),
]

for path, caption, stitle in figures_2d:
    add_image_slide(path, caption, stitle)


# ═══════════════════════════════════════════════════════════
# SLIDES 10-12: Figures 7-9 (3D extruded maps) - JP versions
# ═══════════════════════════════════════════════════════════

figures_3d = [
    ('/home/ubuntu/3d_extruded/3D_ratio_by_anes_v2_jp.png',
     '図7. 3次元押出マップ：色=L003/L008比（硬膜外併用率）、高さ=麻酔科医数。',
     '図7'),
    ('/home/ubuntu/3d_extruded/3D_L008_by_anes_v2_jp.png',
     '図8. 3次元押出マップ：色=L008 SCR（全身麻酔量）、高さ=麻酔科医数。'
     '赤く高い領域は大都市圏の大学病院所在圏。',
     '図8'),
    ('/home/ubuntu/3d_extruded/3D_ratio_by_surgery_v2_jp.png',
     '図9. 3次元押出マップ：色=L003/L008比、高さ=手術あたりGA率。'
     '高さの変動は小さく、手術あたりGA率は比較的均一。色の変動が硬膜外併用率の差を表現。',
     '図9'),
]

for path, caption, stitle in figures_3d:
    add_image_slide(path, caption, stitle)


# ═══════════════════════════════════════════════════════════
# SLIDE 13: STROBE Checklist
# ═══════════════════════════════════════════════════════════

strobe_items = [
    ['1a', '題名と抄録', '表紙および構造化抄録'],
    ['1b', '情報を含む抄録', '全要素を含む構造化抄録'],
    ['2', '背景/根拠', '緒言 第1-3段落'],
    ['3', '目的', '緒言 第5段落'],
    ['4', '研究デザイン', '方法：研究デザインとデータソース'],
    ['5', '対象', '方法：335圏、2022年度データ'],
    ['6a', '参加者/適格基準', '該当なし（生態学的研究）'],
    ['7', '変数', '方法：麻酔コード（表1）'],
    ['8', 'データソース', '方法：研究デザインとデータソース'],
    ['9', 'バイアス', '方法：感度分析、考察：限界'],
    ['10', '研究規模', '方法：335圏、47都道府県'],
    ['11', '量的変数', '方法：統計解析'],
    ['12', '統計手法', '方法：統計解析'],
    ['13', '参加者', '結果：L008データを持つ334圏'],
    ['14', '記述データ', '結果：表2'],
    ['15', 'アウトカムデータ', '結果：地域差セクション'],
    ['16', '主要結果', '結果：感度分析1-3'],
    ['17', 'その他の解析', '結果：合計SCR、ペインクリニック波及'],
    ['18', '主要な知見', '考察：主要所見'],
    ['19', '限界', '考察：長所と限界'],
    ['20', '解釈', '考察：既存文献との比較'],
    ['21', '一般化可能性', '考察：結論'],
    ['22', '資金', '[著者が記入]'],
]

add_table_slide(
    'STROBE チェックリスト（横断研究）',
    ['項目', 'STROBE推奨事項', '本論文での該当箇所'],
    strobe_items,
    col_widths=[1.0, 4.5, 5.5]
)


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════

output_path = '/home/ubuntu/regional_anaesthesia_figures_JP.pptx'
prs.save(output_path)
print(f"Saved: {output_path}")
