#!/usr/bin/env python3
"""
Create EN PPTX with 1 figure/table per slide.
- Code-generated maps: embedded as images (as-is)
- Tables: editable pptx tables
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


def add_title_slide(title_text, subtitle_text=''):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.name = 'Arial'
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER
    return slide


def add_image_slide(image_path, caption, slide_title=''):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Slide title at top
    if slide_title:
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0, 0, 0)

    # Image centered
    if os.path.exists(image_path):
        from PIL import Image
        img = Image.open(image_path)
        img_w, img_h = img.size
        aspect = img_w / img_h

        # Available area: width=12.7in, height=5.8in (leave room for title+caption)
        max_w = Inches(12.0)
        max_h = Inches(5.6)

        if aspect > (12.0 / 5.6):
            w = max_w
            h = int(w / aspect)
        else:
            h = max_h
            w = int(h * aspect)

        # Center horizontally
        left = int((prs.slide_width - w) / 2)
        top = Inches(0.75)

        slide.shapes.add_picture(image_path, left, top, w, h)
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f'[Image not found: {image_path}]'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 0, 0)

    # Caption at bottom
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_table_slide(title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(12.7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0, 0, 0)

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    # Calculate table dimensions
    if col_widths is None:
        total_w = Inches(11.0)
        col_w = total_w / n_cols
        col_widths_emu = [int(col_w)] * n_cols
    else:
        col_widths_emu = [Inches(w) for w in col_widths]

    table_width = sum(col_widths_emu)
    table_left = int((prs.slide_width - table_width) / 2)
    table_top = Inches(1.2)
    row_height = Inches(0.4)
    table_height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths_emu):
        table.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.name = 'Arial'
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(51, 63, 80)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Arial'
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Alternate row colours
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 244, 248)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    return slide


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════

add_title_slide(
    "Regional variation in anaesthesia practice across\n"
    "335 secondary medical areas in Japan",
    "A cross-sectional analysis of national claims data"
)


# ═══════════════════════════════════════════════════════════
# SLIDE 2: Table 1 - Anaesthesia codes
# ═══════════════════════════════════════════════════════════

add_table_slide(
    'Table 1. Anaesthesia procedure codes analysed',
    ['Code', 'Procedure', 'Clinical significance', 'SMAs with data'],
    [
        ['L008', 'Closed-circuit general anaesthesia', 'Primary indicator of GA volume', '334'],
        ['L002', 'Epidural anaesthesia', 'Regional technique (alone or combined)', '307'],
        ['L003', 'Continuous epidural infusion', 'Direct indicator of GA+epidural combination', '331'],
        ['L004', 'Spinal anaesthesia', 'Alternative to GA for suitable procedures', '334'],
        ['L009', 'Anaesthesia management fee I', 'Proxy for specialist anaesthesiologist staffing', '314'],
        ['L100', 'Nerve block (inpatient)', 'Indicator of pain clinic activity', '335'],
    ],
    col_widths=[1.2, 3.5, 4.5, 1.8]
)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: Table 2 - Descriptive statistics
# ═══════════════════════════════════════════════════════════

add_table_slide(
    'Table 2. Distribution of standardised claim ratios (national average = 100)',
    ['Code', 'Procedure', 'P10', 'P50', 'P90', 'CV (%)', 'n'],
    [
        ['L008', 'General anaesthesia', '32.3', '73.2', '131.5', '54.6', '334'],
        ['L002', 'Epidural anaesthesia', '7.6', '73.9', '184.0', '83.2', '307'],
        ['L003', 'Continuous epidural infusion', '19.0', '73.1', '150.0', '64.9', '331'],
        ['L004', 'Spinal anaesthesia', '31.1', '84.3', '169.0', '56.3', '334'],
        ['L009', 'Anaesthesia management fee I', '32.5', '79.1', '149.1', '57.2', '314'],
        ['L100', 'Nerve block (inpatient)', '23.4', '72.2', '174.7', '72.1', '335'],
    ],
    col_widths=[1.0, 3.0, 1.0, 1.0, 1.0, 1.0, 1.0]
)


# ═══════════════════════════════════════════════════════════
# SLIDES 4-9: Figures 1-6 (2D maps)
# ═══════════════════════════════════════════════════════════

figures_2d = [
    ('/home/ubuntu/en_map_L008_scr.png',
     'Figure 1. General anaesthesia (L008) standardised claim ratio by secondary medical area. '
     'Solid lines = prefectural boundaries; dashed lines = SMA boundaries. '
     'Northern Territories shown in white (no SMA designated).',
     'Figure 1'),
    ('/home/ubuntu/en_map_univ_presence.png',
     'Figure 2. University hospital presence by secondary medical area. Red circles indicate '
     'SMAs containing one or more university hospitals (64 of 335 SMAs).',
     'Figure 2'),
    ('/home/ubuntu/en_map_L004_scr.png',
     'Figure 3. Spinal anaesthesia (L004) standardised claim ratio by secondary medical area.',
     'Figure 3'),
    ('/home/ubuntu/en_map_L008_L004_combined.png',
     'Figure 4. Combined general + spinal anaesthesia (L008+L004) standardised claim ratio. '
     'Combining these codes neutralises audit-driven reclassification.',
     'Figure 4'),
    ('/home/ubuntu/en_map_L003_scr.png',
     'Figure 5. Continuous epidural infusion (L003) SCR \u2014 direct indicator of combined '
     'GA\u2013epidural technique. CV=64.9%.',
     'Figure 5'),
    ('/home/ubuntu/en_map_L003_L008_ratio_corrected.png',
     'Figure 6. L003/L008 ratio (combined epidural rate per GA case), adjusted for surgery volume.',
     'Figure 6'),
]

for path, caption, stitle in figures_2d:
    add_image_slide(path, caption, stitle)


# ═══════════════════════════════════════════════════════════
# SLIDES 10-12: Figures 7-9 (3D extruded maps)
# ═══════════════════════════════════════════════════════════

figures_3d = [
    ('/home/ubuntu/3d_extruded/3D_ratio_by_anes_v2.png',
     'Figure 7. 3D extruded map: colour = L003/L008 ratio (combined epidural rate), '
     'height = anaesthesiologist count.',
     'Figure 7'),
    ('/home/ubuntu/3d_extruded/3D_L008_by_anes_v2.png',
     'Figure 8. 3D extruded map: colour = L008 SCR (general anaesthesia volume), '
     'height = anaesthesiologist count.',
     'Figure 8'),
    ('/home/ubuntu/3d_extruded/3D_ratio_by_surgery_v2.png',
     'Figure 9. 3D extruded map: colour = L003/L008 ratio, '
     'height = GA per surgery ratio (L008/K-chapter SCR).',
     'Figure 9'),
]

for path, caption, stitle in figures_3d:
    add_image_slide(path, caption, stitle)


# ═══════════════════════════════════════════════════════════
# SLIDE 13: STROBE Checklist
# ═══════════════════════════════════════════════════════════

strobe_items = [
    ['1a', 'Title and abstract', 'Title page and structured abstract'],
    ['1b', 'Informative abstract', 'Structured abstract with all required elements'],
    ['2', 'Background/rationale', 'Introduction paragraphs 1\u20133'],
    ['3', 'Objectives', 'Introduction paragraph 5'],
    ['4', 'Study design', 'Methods: Study design and data sources'],
    ['5', 'Setting', 'Methods: 335 SMAs, FY2022 data'],
    ['6a', 'Participants/eligibility', 'N/A (ecological study)'],
    ['7', 'Variables', 'Methods: Anaesthesia codes (Table 1)'],
    ['8', 'Data sources', 'Methods: Study design and data sources'],
    ['9', 'Bias', 'Methods: Sensitivity analyses; Discussion: Limitations'],
    ['10', 'Study size', 'Methods: 335 SMAs, 47 prefectures'],
    ['11', 'Quantitative variables', 'Methods: Statistical analysis'],
    ['12', 'Statistical methods', 'Methods: Statistical analysis'],
    ['13', 'Participants', 'Results: 334 SMAs with L008 data'],
    ['14', 'Descriptive data', 'Results: Table 2'],
    ['15', 'Outcome data', 'Results: Regional variation section'],
    ['16', 'Main results', 'Results: Sensitivity analyses 1\u20133'],
    ['17', 'Other analyses', 'Results: Combined SCR, Pain clinic spillover'],
    ['18', 'Key results', 'Discussion: Principal findings'],
    ['19', 'Limitations', 'Discussion: Strengths and limitations'],
    ['20', 'Interpretation', 'Discussion: Comparison with literature'],
    ['21', 'Generalisability', 'Discussion: Conclusions'],
    ['22', 'Funding', '[To be completed]'],
]

add_table_slide(
    'STROBE Checklist (Cross-sectional study)',
    ['Item', 'STROBE recommendation', 'Location in manuscript'],
    strobe_items,
    col_widths=[1.0, 4.5, 5.5]
)


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════

output_path = '/home/ubuntu/regional_anaesthesia_figures_EN.pptx'
prs.save(output_path)
print(f"Saved: {output_path}")
